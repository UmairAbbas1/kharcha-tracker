"""
generate_pptx.py
Generates kharcha-tracker-phase4.pptx — a professional 10-slide PowerPoint
for Kharcha Tracker, Pakistan's AI-Native Expense Tracker SaaS.

Usage:
    pip install python-pptx
    python generate_pptx.py

If python-pptx is not installed, the script automatically falls back to
generate_pptx_node.js (requires Node.js + pptxgenjs).
"""

import os, subprocess, sys

# ── Guard: try python-pptx, else delegate to Node.js companion ────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    node_script = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "generate_pptx_node.js")
    print("python-pptx not installed — delegating to Node.js companion …")
    sys.exit(subprocess.run(["node", node_script]).returncode)

# ── Colour palette ────────────────────────────────────────────────────────────
ROYAL_BLUE = RGBColor(65, 105, 225)
WHITE      = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(238, 242, 255)
DARK_TEXT  = RGBColor(30, 30, 50)
MID_GREY   = RGBColor(180, 180, 190)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT    = "Calibri"


# ── Low-level helpers ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, font_size, bold=False,
                color=None, align=PP_ALIGN.LEFT, italic=False):
    color = color or DARK_TEXT
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = FONT
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def set_cell_bg(cell, rgb):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(rgb.r, rgb.g, rgb.b))


def style_cell(cell, text, font_size=12, bold=False, fg=None, bg=None,
               align=PP_ALIGN.LEFT):
    fg = fg or DARK_TEXT
    if bg:
        set_cell_bg(cell, bg)
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name  = FONT
            run.font.size  = Pt(font_size)
            run.font.bold  = bold
            run.font.color.rgb = fg


def blue_title(slide, text, top=Inches(0.22)):
    add_textbox(slide, Inches(0.5), top, Inches(12.33), Inches(0.7),
                text, 30, bold=True, color=ROYAL_BLUE)


def top_bottom_bars(slide):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_rgb=ROYAL_BLUE)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_rgb=ROYAL_BLUE)


def divider(slide, top=Inches(1.05)):
    add_rect(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), fill_rgb=ROYAL_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def slide1(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=ROYAL_BLUE)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), fill_rgb=WHITE)
    add_rect(s, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_rgb=WHITE)

    add_textbox(s, Inches(1), Inches(1.7), Inches(11.33), Inches(1.3),
                "Kharcha Tracker", 54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.15), Inches(11.33), Inches(0.75),
                "Pakistan\u2019s First AI-Native Expense Tracker",
                24, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(4.5), Inches(4.2), Inches(4.33), Inches(0.04), fill_rgb=WHITE)
    add_textbox(s, Inches(1), Inches(4.4), Inches(11.33), Inches(0.6),
                "Phase 4 R\u0026D Plan \u2014 July 2026", 18, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
def slide2(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "The Problem")
    divider(s, Inches(1.05))

    bullets = [
        "Pakistanis track expenses manually \u2014 in notebooks or WhatsApp messages",
        "Global apps (Toshl, Spendee) don\u2019t understand PKR, Easypaisa, or JazzCash",
        "No app sends WhatsApp budget alerts or parses bank SMS automatically",
    ]
    top = Inches(1.3)
    for b in bullets:
        add_textbox(s, Inches(0.5), top, Inches(12.33), Inches(1.5),
                    "\u25C6  " + b, 20)
        top += Inches(1.6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — USP Table
# ══════════════════════════════════════════════════════════════════════════════
def slide3(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "Our Unique Advantage")
    divider(s, Inches(1.0))

    rows_data = [
        ("Competitors",            "Kharcha Tracker"),
        ("AI OCR is paid-only",    "AI Receipt Scanning \u2014 FREE"),
        ("No WhatsApp alerts",     "WhatsApp Budget Alerts"),
        ("USD-first, global design", "PKR-native, Pakistani receipts"),
        ("No voice entry",         "Urdu + English Voice Entry (Phase 4)"),
        ("No bank SMS parsing",    "HBL/MCB/UBL SMS Auto-Parse (Phase 4)"),
        ("Paid team features",     "Free multi-user workspaces"),
    ]

    tbl = s.shapes.add_table(
        len(rows_data), 2, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.9)
    ).table
    tbl.columns[0].width = Inches(5.5)
    tbl.columns[1].width = Inches(6.83)

    for r, (c1, c2) in enumerate(rows_data):
        is_hdr = r == 0
        is_odd = r % 2 == 1
        bg = ROYAL_BLUE if is_hdr else (LIGHT_BLUE if is_odd else WHITE)
        fg = WHITE if is_hdr else DARK_TEXT
        for ci, txt in enumerate([c1, c2]):
            cell = tbl.cell(r, ci)
            style_cell(cell, txt, font_size=13, bold=is_hdr, fg=fg, bg=bg,
                       align=PP_ALIGN.CENTER if is_hdr else PP_ALIGN.LEFT)
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.05)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What We Built
# ══════════════════════════════════════════════════════════════════════════════
def slide4(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "What\u2019s Already Built")
    divider(s, Inches(1.05))

    items = [
        ("Phase 1 \u2014 Multi-tenant SaaS",   "Supabase auth, workspaces, roles, RLS security"),
        ("Phase 2 \u2014 Smart Budget Alerts", "Email + WhatsApp notifications at 80 / 90 / 100% spend"),
        ("Phase 3 \u2014 AI Receipt Scanning", "Groq Llama 4 Scout OCR, mobile UI, optimistic updates"),
    ]
    top = Inches(1.25)
    for label, detail in items:
        add_rect(s, Inches(0.5), top, Inches(12.33), Inches(1.5), fill_rgb=LIGHT_BLUE)
        add_textbox(s, Inches(0.7), top + Inches(0.1), Inches(11.8), Inches(0.55),
                    "\u2705  " + label, 18, bold=True, color=ROYAL_BLUE)
        add_textbox(s, Inches(1.1), top + Inches(0.65), Inches(11.3), Inches(0.65),
                    detail, 14)
        top += Inches(1.65)

    add_textbox(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.4),
                "All built on zero-cost infrastructure",
                12, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Phase 4 Features Table
# ══════════════════════════════════════════════════════════════════════════════
def slide5(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "Phase 4 \u2014 What We Build Next")
    divider(s, Inches(1.0))

    rows_data = [
        ("Feature",                     "Effort",    "Why It Wins"),
        ("\U0001f3a4  Voice Expense Entry",  "1 day",     "Say \u2018Rs 450 biryani\u2019 \u2014 done. No competitor offers this"),
        ("\U0001f4f1  Bank SMS Parsing",     "2 days",    "HBL/MCB/UBL SMS auto-creates expense drafts"),
        ("\U0001f4ca  Monthly AI Summary",   "1 day",     "Plain-language spending report delivered by email"),
        ("\U0001f4e4  CSV / Excel Export",   "2 hours",   "Every user wants data portability"),
        ("\U0001f3a8  Google Stitch UI",     "1 day",     "AI-generated polished UI via Google\u2019s design tool"),
        ("\U0001f501  AI Provider Failover", "2 hours",   "Groq + Gemini fallback \u2014 never breaks in production"),
        ("Total Phase 4",               "\u22486 days", "Full production-ready SaaS"),
    ]

    tbl = s.shapes.add_table(
        len(rows_data), 3, Inches(0.5), Inches(1.1), Inches(12.33), Inches(6.1)
    ).table
    for i, w in enumerate([Inches(4.5), Inches(1.8), Inches(6.03)]):
        tbl.columns[i].width = w

    for r, row in enumerate(rows_data):
        is_hdr   = r == 0
        is_total = r == len(rows_data) - 1
        is_odd   = r % 2 == 1
        bg = ROYAL_BLUE if (is_hdr or is_total) else (LIGHT_BLUE if is_odd else WHITE)
        fg = WHITE if (is_hdr or is_total) else DARK_TEXT
        fs = 13 if is_hdr else 12
        for ci, txt in enumerate(row):
            cell = tbl.cell(r, ci)
            al = PP_ALIGN.CENTER if (is_hdr or ci == 1) else PP_ALIGN.LEFT
            style_cell(cell, txt, font_size=fs, bold=(is_hdr or is_total), fg=fg, bg=bg, align=al)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Google Stitch
# ══════════════════════════════════════════════════════════════════════════════
def slide6(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "Google Stitch \u2014 AI UI Design")
    divider(s, Inches(1.05))

    add_textbox(s, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.85),
                ("Google\u2019s new AI design tool (launched I/O 2025) generates complete "
                 "UI screens from a text prompt in under 60 seconds"),
                17, italic=True, color=ROYAL_BLUE)

    bullets = [
        "Exports clean HTML/CSS and Figma files \u2014 no manual design work",
        "Connects to Kiro AI via Model Context Protocol (MCP) \u2014 designs flow straight into code",
        "Prompt used: \u2018Pakistani expense tracker, royal blue #4169E1, mobile-first, glassmorphism\u2019",
    ]
    top = Inches(2.25)
    for b in bullets:
        add_textbox(s, Inches(0.6), top, Inches(12.1), Inches(1.0), "\u2192  " + b, 18)
        top += Inches(1.1)

    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.62), fill_rgb=LIGHT_BLUE)
    add_textbox(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.5),
                "stitch.withgoogle.com  \u2014  Free", 14, bold=True,
                color=ROYAL_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
def slide7(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "Technology Stack \u2014 All Free Tier")
    divider(s, Inches(1.05))

    left  = ["Supabase (Database + Auth)", "React 18 + Vite + Tailwind CSS",
             "Node.js + Express", "Groq AI (OCR + Text)", "Resend (Email alerts)"]
    right = ["Groq Whisper (Voice)", "Google Gemini (AI fallback)",
             "Google Stitch (UI design)", "Vercel + Render (Deployment)",
             "cron-job.org (Scheduled jobs)"]

    for panel_l, label, items in [
        (Inches(0.5),  "Already in use",   left),
        (Inches(7.0),  "Adding in Phase 4", right),
    ]:
        add_rect(s, panel_l, Inches(1.2), Inches(5.8), Inches(5.85), fill_rgb=LIGHT_BLUE)
        add_textbox(s, panel_l + Inches(0.15), Inches(1.28), Inches(5.5), Inches(0.5),
                    label, 15, bold=True, color=ROYAL_BLUE)
        add_rect(s, panel_l + Inches(0.15), Inches(1.82), Inches(5.5), Inches(0.04), fill_rgb=ROYAL_BLUE)
        top = Inches(1.98)
        for item in items:
            add_textbox(s, panel_l + Inches(0.3), top, Inches(5.3), Inches(0.7),
                        "\u2714  " + item, 15)
            top += Inches(0.82)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Pricing
# ══════════════════════════════════════════════════════════════════════════════
def slide8(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "Business Model")
    divider(s, Inches(1.05))

    cards = [
        ("Free",  "Rs 0 / month",     WHITE,      ROYAL_BLUE, "1 workspace, manual entry, basic alerts"),
        ("Pro",   "Rs 499 / month",   ROYAL_BLUE, WHITE,      "AI OCR, voice entry, CSV export, monthly summary"),
        ("Team",  "Rs 1,499 / month", LIGHT_BLUE, ROYAL_BLUE, "Multi-user, SMS parsing, WhatsApp alerts"),
    ]
    card_w = Inches(3.9)
    x = Inches(0.5)
    for tier, price, bg, fg, detail in cards:
        add_rect(s, x, Inches(1.28), card_w, Inches(5.2), fill_rgb=bg,
                 line_rgb=ROYAL_BLUE, line_width=Pt(1.5))
        add_textbox(s, x, Inches(1.5),  card_w, Inches(0.65),
                    tier,  26, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.3), Inches(2.22), card_w - Inches(0.6), Inches(0.04), fill_rgb=fg)
        add_textbox(s, x, Inches(2.35), card_w, Inches(0.75),
                    price, 22, bold=True, color=fg, align=PP_ALIGN.CENTER)

        detail_box = s.shapes.add_textbox(x + Inches(0.2), Inches(3.25),
                                           card_w - Inches(0.4), Inches(2.9))
        detail_box.word_wrap = True
        tf = detail_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = detail
        run.font.name = FONT
        run.font.size = Pt(14)
        run.font.color.rgb = fg

        x += card_w + Inches(0.27)

    add_textbox(s, Inches(0.5), Inches(6.72), Inches(12.33), Inches(0.4),
                "Pakistan-first pricing.  Approx $1.80 / $5.40 USD",
                12, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ══════════════════════════════════════════════════════════════════════════════
def slide9(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    top_bottom_bars(s)
    blue_title(s, "6-Day Implementation Plan")
    divider(s, Inches(1.05))

    timeline = [
        ("Day 1",         "CSV Export + OCR Provider Abstraction"),
        ("Day 2\u20133",  "Bank SMS Parsing"),
        ("Day 4",         "Voice Expense Entry"),
        ("Day 5",         "Monthly AI Summary Email"),
        ("Day 6",         "Google Stitch UI Refresh + Deploy to Vercel + Render"),
    ]
    top = Inches(1.28)
    for day, task in timeline:
        add_rect(s, Inches(0.5), top, Inches(1.6), Inches(0.9), fill_rgb=ROYAL_BLUE)
        add_textbox(s, Inches(0.5), top + Inches(0.1), Inches(1.6), Inches(0.7),
                    day, 15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(2.25), top, Inches(10.6), Inches(0.9), fill_rgb=LIGHT_BLUE)
        add_textbox(s, Inches(2.45), top + Inches(0.1), Inches(10.2), Inches(0.7), task, 16)
        add_rect(s, Inches(2.1), top + Inches(0.38), Inches(0.15), Inches(0.15), fill_rgb=ROYAL_BLUE)
        top += Inches(1.05)

    add_textbox(s, Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.45),
                "Built entirely with Kiro AI on zero-cost infrastructure",
                13, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Closing
# ══════════════════════════════════════════════════════════════════════════════
def slide10(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=ROYAL_BLUE)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), fill_rgb=WHITE)
    add_rect(s, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_rgb=WHITE)

    add_textbox(s, Inches(1), Inches(1.7), Inches(11.33), Inches(1.3),
                "Kharcha Tracker", 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.1), Inches(11.33), Inches(0.75),
                "Pakistan\u2019s smartest expense tracker", 22, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(4.5), Inches(4.0), Inches(4.33), Inches(0.04), fill_rgb=WHITE)
    add_textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.65),
                "Built by Umair Abbas  \u00B7  AI Engineering Intern  \u00B7  July 2026",
                16, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(4.95), Inches(11.33), Inches(0.5),
                "github.com/UmairAbbas1/kharcha-tracker",
                13, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    print("Building slides …")
    slide1(prs);  print("  \u2713 Slide  1 \u2014 Title")
    slide2(prs);  print("  \u2713 Slide  2 \u2014 The Problem")
    slide3(prs);  print("  \u2713 Slide  3 \u2014 USP table")
    slide4(prs);  print("  \u2713 Slide  4 \u2014 What We Built")
    slide5(prs);  print("  \u2713 Slide  5 \u2014 Phase 4 Features")
    slide6(prs);  print("  \u2713 Slide  6 \u2014 Google Stitch")
    slide7(prs);  print("  \u2713 Slide  7 \u2014 Tech Stack")
    slide8(prs);  print("  \u2713 Slide  8 \u2014 Pricing")
    slide9(prs);  print("  \u2713 Slide  9 \u2014 Timeline")
    slide10(prs); print("  \u2713 Slide 10 \u2014 Closing")

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "kharcha-tracker-phase4.pptx")
    prs.save(out_path)
    print(f"\n\u2705  Saved \u2192 {out_path}")


if __name__ == "__main__":
    main()
